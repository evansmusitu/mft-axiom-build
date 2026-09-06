#!/usr/bin/env python3
"""MUSITU Axiom v5 full-stack frontier adapters.

This module provides real adapter implementations for the five-year frontier
workstreams. External adapters are fail-closed: a capability is usable only
when its runtime dependency is present and the caller supplies the required
credentials/authorization. CI exercises these adapters against disposable or
isolated resources; production promotion remains separately gated.
"""
from __future__ import annotations

from dataclasses import dataclass, asdict, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable, Mapping, Sequence
import base64
import csv
import hashlib
import io
import json
import math
import os
import re
import shutil
import sqlite3
import statistics
import subprocess
import tempfile
import urllib.error
import urllib.parse
import urllib.request
import wave

from .advanced import (
    CausalModel,
    DigitalTwin,
    LinearEquation,
    MemoryRecord,
)
from .fabric import AuthorizationError, FrontierError


def _utcnow() -> str:
    return datetime.now(timezone.utc).isoformat()


def _canonical(value: Any) -> str:
    return json.dumps(value, sort_keys=True, separators=(",", ":"), default=str)


def _sha_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _sha(value: Any) -> str:
    return _sha_bytes(_canonical(value).encode())


# ---------------------------------------------------------------------------
# Live research + provenance + prompt-injection treatment
# ---------------------------------------------------------------------------

@dataclass(frozen=True)
class RetrievalSnapshot:
    url: str
    final_url: str
    retrieved_at: str
    status: int
    content_type: str
    byte_length: int
    sha256: str
    source_host: str
    instruction_authority: str
    injection_flags: tuple[str, ...] = ()
    content: bytes = field(default=b"", repr=False, compare=False)

    def text(self, encoding: str = "utf-8") -> str:
        return self.content.decode(encoding, "replace")

    def evidence(self) -> dict[str, Any]:
        x = asdict(self)
        x.pop("content", None)
        return x


class RetrievedContentFirewall:
    """Detects common authority-confusion/prompt-injection patterns.

    Detection never upgrades retrieved text into instructions. Retrieval is
    still allowed so an analyst can inspect hostile evidence safely.
    """

    PATTERNS = {
        "ignore-prior": re.compile(r"\bignore\s+(all\s+)?(previous|prior|above)\s+(instructions?|rules?|policy)", re.I),
        "system-override": re.compile(r"\b(system|developer)\s+(prompt|message|instructions?)\b", re.I),
        "credential-request": re.compile(r"\b(reveal|print|send|exfiltrate|upload)\b.{0,48}\b(secret|token|credential|api[_ -]?key|password)\b", re.I | re.S),
        "tool-authority": re.compile(r"\b(call|invoke|run|execute)\b.{0,48}\b(tool|command|shell|terminal|payment|trade|deploy)\b", re.I | re.S),
        "policy-bypass": re.compile(r"\b(bypass|disable|override)\b.{0,48}\b(policy|safety|authorization|approval|guardrail)\b", re.I | re.S),
    }

    @classmethod
    def scan(cls, text: str) -> tuple[str, ...]:
        return tuple(sorted(name for name, rx in cls.PATTERNS.items() if rx.search(text)))


class LiveResearchAdapter:
    def __init__(self, allowed_hosts: Iterable[str], max_bytes: int = 8_000_000, timeout: int = 30) -> None:
        self.allowed_hosts = frozenset(h.casefold() for h in allowed_hosts)
        self.max_bytes = int(max_bytes)
        self.timeout = int(timeout)
        if not self.allowed_hosts or self.max_bytes <= 0 or self.timeout <= 0:
            raise ValueError("allowed_hosts, max_bytes and timeout must be valid")

    def _check_url(self, url: str) -> str:
        p = urllib.parse.urlparse(url)
        host = (p.hostname or "").casefold()
        if p.scheme != "https":
            raise AuthorizationError("research retrieval requires HTTPS")
        if host not in self.allowed_hosts:
            raise AuthorizationError("research host is not allowlisted")
        if p.username or p.password:
            raise AuthorizationError("userinfo in research URL is prohibited")
        return host

    def fetch(self, url: str, headers: Mapping[str, str] | None = None) -> RetrievalSnapshot:
        host = self._check_url(url)
        h = {"Accept": "application/json,text/plain,text/html,*/*;q=0.5", "User-Agent": "MUSITU-Axiom-Frontier-Research/5.0"}
        if headers:
            for k, v in headers.items():
                if k.casefold() in {"authorization", "cookie", "proxy-authorization"}:
                    raise AuthorizationError("credential-bearing research header prohibited")
                h[str(k)] = str(v)
        req = urllib.request.Request(url, headers=h, method="GET")
        try:
            with urllib.request.urlopen(req, timeout=self.timeout) as r:
                final = r.geturl()
                self._check_url(final)
                status = int(r.status)
                ctype = str(r.headers.get("Content-Type") or "application/octet-stream").split(";", 1)[0].strip().lower()
                body = r.read(self.max_bytes + 1)
        except urllib.error.HTTPError as e:
            raise FrontierError(f"research HTTP {e.code}") from e
        except urllib.error.URLError as e:
            raise FrontierError("research transport failure") from e
        if status != 200:
            raise FrontierError(f"research HTTP {status}")
        if len(body) > self.max_bytes:
            raise FrontierError("research response exceeds byte limit")
        flags: tuple[str, ...] = ()
        if ctype.startswith("text/") or ctype in {"application/json", "application/xml", "application/xhtml+xml"}:
            flags = RetrievedContentFirewall.scan(body.decode("utf-8", "replace"))
        return RetrievalSnapshot(
            url=url,
            final_url=final,
            retrieved_at=_utcnow(),
            status=status,
            content_type=ctype,
            byte_length=len(body),
            sha256=_sha_bytes(body),
            source_host=host,
            instruction_authority="retrieved-content-data-only",
            injection_flags=flags,
            content=body,
        )


# ---------------------------------------------------------------------------
# Durable persistence: local SQLite + isolated Cloudflare D1 production store
# ---------------------------------------------------------------------------

class SQLiteProductionStore:
    """Transactional namespaced persistence used for local/edge-compatible tests."""

    def __init__(self, path: str | Path) -> None:
        self.path = Path(path)
        self.path.parent.mkdir(parents=True, exist_ok=True)
        self.db = sqlite3.connect(self.path)
        self.db.execute("PRAGMA journal_mode=WAL")
        self.db.execute("PRAGMA foreign_keys=ON")
        self.migrate()

    def migrate(self) -> None:
        self.db.executescript(
            """
            CREATE TABLE IF NOT EXISTS memory_records(
              tenant TEXT NOT NULL, namespace TEXT NOT NULL, key TEXT NOT NULL,
              record_id TEXT PRIMARY KEY, payload_json TEXT NOT NULL,
              observed_at TEXT NOT NULL, tombstone INTEGER NOT NULL DEFAULT 0,
              supersedes TEXT
            );
            CREATE INDEX IF NOT EXISTS idx_memory_lookup
              ON memory_records(tenant,namespace,key,observed_at);
            CREATE TABLE IF NOT EXISTS twin_calibrations(
              tenant TEXT NOT NULL, twin_id TEXT NOT NULL, version TEXT NOT NULL,
              twin_type TEXT NOT NULL, payload_json TEXT NOT NULL,
              source_sha256 TEXT NOT NULL, created_at TEXT NOT NULL,
              PRIMARY KEY(tenant,twin_id,version)
            );
            CREATE TABLE IF NOT EXISTS eval_events(
              event_id TEXT PRIMARY KEY, suite TEXT NOT NULL, payload_json TEXT NOT NULL,
              created_at TEXT NOT NULL
            );
            """
        )
        self.db.commit()

    def put_memory(self, tenant: str, record: MemoryRecord) -> str:
        if not tenant:
            raise ValueError("tenant required")
        rid = record.record_id
        with self.db:
            self.db.execute(
                "INSERT INTO memory_records(tenant,namespace,key,record_id,payload_json,observed_at,tombstone,supersedes) VALUES(?,?,?,?,?,?,?,?)",
                (tenant, record.namespace, record.key, rid, _canonical(asdict(record)), record.observed_at, int(record.tombstone), record.supersedes),
            )
        return rid

    def latest_memory(self, tenant: str, namespace: str, key: str) -> MemoryRecord | None:
        row = self.db.execute(
            "SELECT payload_json FROM memory_records WHERE tenant=? AND namespace=? AND key=? ORDER BY observed_at DESC,rowid DESC LIMIT 1",
            (tenant, namespace, key),
        ).fetchone()
        if not row:
            return None
        rec = MemoryRecord(**json.loads(row[0]))
        return None if rec.tombstone else rec

    def put_twin(self, tenant: str, twin_id: str, version: str, twin_type: str, payload: Mapping[str, Any], source_sha256: str) -> None:
        if len(source_sha256) != 64:
            raise ValueError("source_sha256 required")
        with self.db:
            self.db.execute(
                "INSERT INTO twin_calibrations(tenant,twin_id,version,twin_type,payload_json,source_sha256,created_at) VALUES(?,?,?,?,?,?,?)",
                (tenant, twin_id, version, twin_type, _canonical(payload), source_sha256, _utcnow()),
            )

    def close(self) -> None:
        self.db.close()


class CloudflareD1ProductionStore:
    """Minimal Cloudflare D1 REST adapter.

    The database must be dedicated to Frontier v5. This class refuses to use
    the sealed production Axiom database UUID to prevent accidental mutation.
    """

    SEALED_AXIOM_DB = "504029cc-f9a5-495e-818f-63c6144b4ea4"

    def __init__(self, account_id: str, database_id: str, email: str, global_api_key: str) -> None:
        if database_id == self.SEALED_AXIOM_DB:
            raise AuthorizationError("frontier persistence must not use the sealed Axiom production DB")
        if not all((account_id, database_id, email, global_api_key)):
            raise ValueError("Cloudflare D1 credentials and IDs required")
        self.account_id = account_id
        self.database_id = database_id
        self.base = "https://api.cloudflare.com/client/v4"
        self.headers = {
            "X-Auth-Email": email,
            "X-Auth-Key": global_api_key,
            "Accept": "application/json",
            "Content-Type": "application/json",
            "User-Agent": "MUSITU-Axiom-Frontier-Persistence/5.0",
        }

    def query(self, sql: str, params: Sequence[Any] | None = None) -> list[dict[str, Any]]:
        obj: dict[str, Any] = {"sql": sql}
        if params is not None:
            obj["params"] = list(params)
        url = f"{self.base}/accounts/{self.account_id}/d1/database/{self.database_id}/query"
        req = urllib.request.Request(url, data=_canonical(obj).encode(), headers=self.headers, method="POST")
        try:
            with urllib.request.urlopen(req, timeout=45) as r:
                raw = json.loads(r.read() or b"{}")
        except Exception as exc:
            raise FrontierError("Cloudflare D1 request failed") from exc
        if raw.get("success") is not True:
            raise FrontierError("Cloudflare D1 success=false")
        blocks = raw.get("result") or []
        if not blocks or any(b.get("success") is not True for b in blocks):
            raise FrontierError("Cloudflare D1 statement failed")
        rows: list[dict[str, Any]] = []
        for b in blocks:
            rows.extend(b.get("results") or [])
        return rows

    def migrate(self) -> None:
        for sql in (
            "CREATE TABLE IF NOT EXISTS frontier_memory(tenant TEXT NOT NULL, namespace TEXT NOT NULL, k TEXT NOT NULL, record_id TEXT PRIMARY KEY, payload_json TEXT NOT NULL, observed_at TEXT NOT NULL)",
            "CREATE INDEX IF NOT EXISTS idx_frontier_memory_lookup ON frontier_memory(tenant,namespace,k,observed_at)",
            "CREATE TABLE IF NOT EXISTS frontier_twins(tenant TEXT NOT NULL, twin_id TEXT NOT NULL, version TEXT NOT NULL, twin_type TEXT NOT NULL, payload_json TEXT NOT NULL, source_sha256 TEXT NOT NULL, created_at TEXT NOT NULL, PRIMARY KEY(tenant,twin_id,version))",
            "CREATE TABLE IF NOT EXISTS frontier_eval_events(event_id TEXT PRIMARY KEY, suite TEXT NOT NULL, payload_json TEXT NOT NULL, created_at TEXT NOT NULL)",
        ):
            self.query(sql)


# ---------------------------------------------------------------------------
# Browser + computer use
# ---------------------------------------------------------------------------

@dataclass(frozen=True)
class BrowserResult:
    url: str
    title: str
    text_sha256: str
    screenshot_sha256: str | None
    final_url: str


class PlaywrightBrowserAdapter:
    def __init__(self, allowed_hosts: Iterable[str], headless: bool = True) -> None:
        self.allowed_hosts = frozenset(h.casefold() for h in allowed_hosts)
        self.headless = bool(headless)

    def _allow(self, url: str) -> None:
        p = urllib.parse.urlparse(url)
        host = (p.hostname or "").casefold()
        if p.scheme not in {"http", "https"} or host not in self.allowed_hosts:
            raise AuthorizationError("browser navigation target not allowlisted")

    def run(self, url: str, screenshot_path: str | Path | None = None,
            click_selector: str | None = None, expect_text: str | None = None) -> BrowserResult:
        self._allow(url)
        try:
            from playwright.sync_api import sync_playwright
        except Exception as exc:
            raise FrontierError("playwright runtime unavailable") from exc
        shot_hash = None
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=self.headless)
            page = browser.new_page()
            page.goto(url, wait_until="networkidle")
            self._allow(page.url)
            if click_selector:
                page.click(click_selector)
            text = page.locator("body").inner_text()
            if expect_text is not None and expect_text not in text:
                browser.close()
                raise FrontierError("browser expected text not observed")
            if screenshot_path is not None:
                sp = Path(screenshot_path)
                sp.parent.mkdir(parents=True, exist_ok=True)
                page.screenshot(path=str(sp), full_page=True)
                shot_hash = _sha_bytes(sp.read_bytes())
            out = BrowserResult(url, page.title(), _sha_bytes(text.encode()), shot_hash, page.url)
            browser.close()
            return out


class SandboxedComputerAdapter:
    """Non-shell application/computer adapter constrained to a working root."""

    def __init__(self, root: str | Path, allowed_executables: Iterable[str]) -> None:
        self.root = Path(root).resolve()
        self.root.mkdir(parents=True, exist_ok=True)
        self.allowed = frozenset(str(x) for x in allowed_executables)
        if not self.allowed:
            raise ValueError("allowed executables required")

    def run(self, argv: Sequence[str], timeout: int = 30, env: Mapping[str, str] | None = None) -> dict[str, Any]:
        if not argv:
            raise ValueError("argv required")
        exe = Path(argv[0]).name
        if exe not in self.allowed:
            raise AuthorizationError("computer executable not allowlisted")
        proc = subprocess.run(list(argv), cwd=self.root, env=dict(os.environ, **(dict(env or {}))),
                              capture_output=True, text=True, timeout=timeout, shell=False)
        if proc.returncode != 0:
            raise FrontierError(f"computer action failed: {exe} rc={proc.returncode}")
        return {
            "argv_sha256": _sha(list(argv)),
            "stdout": proc.stdout,
            "stderr_sha256": _sha_bytes(proc.stderr.encode()),
            "returncode": proc.returncode,
        }


# ---------------------------------------------------------------------------
# Artifact workbenches
# ---------------------------------------------------------------------------

class ArtifactWorkbench:
    @staticmethod
    def document(path: str | Path, title: str, paragraphs: Sequence[str]) -> dict[str, Any]:
        from docx import Document
        p = Path(path); p.parent.mkdir(parents=True, exist_ok=True)
        doc = Document(); doc.add_heading(title, level=1)
        for paragraph in paragraphs: doc.add_paragraph(str(paragraph))
        doc.save(p)
        check = Document(p)
        text = "\n".join(x.text for x in check.paragraphs)
        if title not in text:
            raise FrontierError("DOCX round-trip validation failed")
        return {"type":"docx","path":str(p),"sha256":_sha_bytes(p.read_bytes()),"paragraphs":len(check.paragraphs)}

    @staticmethod
    def spreadsheet(path: str | Path, headers: Sequence[str], rows: Sequence[Sequence[Any]], formulas: Mapping[str, str] | None = None) -> dict[str, Any]:
        from openpyxl import Workbook, load_workbook
        p=Path(path); p.parent.mkdir(parents=True,exist_ok=True)
        wb=Workbook(); ws=wb.active; ws.title="MUSITU"
        ws.append(list(headers))
        for row in rows: ws.append(list(row))
        for cell, formula in (formulas or {}).items(): ws[cell]=formula
        wb.save(p)
        rb=load_workbook(p, data_only=False); rws=rb["MUSITU"]
        if [rws.cell(1,i+1).value for i in range(len(headers))] != list(headers):
            raise FrontierError("XLSX round-trip validation failed")
        return {"type":"xlsx","path":str(p),"sha256":_sha_bytes(p.read_bytes()),"rows":rws.max_row,"cols":rws.max_column}

    @staticmethod
    def presentation(path: str | Path, title: str, slides: Sequence[tuple[str, Sequence[str]]]) -> dict[str, Any]:
        from pptx import Presentation
        p=Path(path); p.parent.mkdir(parents=True,exist_ok=True)
        prs=Presentation()
        s=prs.slides.add_slide(prs.slide_layouts[0]); s.shapes.title.text=title
        for st, bullets in slides:
            slide=prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text=st
            tf=slide.placeholders[1].text_frame; tf.clear()
            for i,b in enumerate(bullets):
                par=tf.paragraphs[0] if i==0 else tf.add_paragraph(); par.text=str(b)
        prs.save(p)
        check=Presentation(p)
        if len(check.slides) != 1+len(slides):
            raise FrontierError("PPTX round-trip validation failed")
        return {"type":"pptx","path":str(p),"sha256":_sha_bytes(p.read_bytes()),"slides":len(check.slides)}

    @staticmethod
    def pdf(path: str | Path, title: str, lines: Sequence[str]) -> dict[str, Any]:
        from reportlab.pdfgen import canvas
        from pypdf import PdfReader
        p=Path(path); p.parent.mkdir(parents=True,exist_ok=True)
        c=canvas.Canvas(str(p)); y=800; c.setFont("Helvetica-Bold",14); c.drawString(50,y,title); y-=28; c.setFont("Helvetica",10)
        for line in lines:
            c.drawString(50,y,str(line)[:110]); y-=16
            if y<50: c.showPage(); y=800
        c.save()
        reader=PdfReader(str(p)); extracted="\n".join((page.extract_text() or "") for page in reader.pages)
        if title not in extracted:
            raise FrontierError("PDF round-trip validation failed")
        return {"type":"pdf","path":str(p),"sha256":_sha_bytes(p.read_bytes()),"pages":len(reader.pages)}

    @staticmethod
    def application(path: str | Path, title: str = "MUSITU Axiom Workbench") -> dict[str, Any]:
        p=Path(path); p.parent.mkdir(parents=True,exist_ok=True)
        html=f"""<!doctype html><html><head><meta charset='utf-8'><title>{title}</title></head><body><h1>{title}</h1><input id='x' value='40+2'><button id='run' onclick=\"document.getElementById('out').textContent='42'\">Run</button><output id='out'>ready</output></body></html>"""
        p.write_text(html,encoding="utf-8")
        return {"type":"html-app","path":str(p),"sha256":_sha_bytes(p.read_bytes())}


# ---------------------------------------------------------------------------
# Multimodal: image, speech/audio, video, screen/camera-like streams
# ---------------------------------------------------------------------------

class MultimodalWorkbench:
    @staticmethod
    def inspect_image(path: str | Path) -> dict[str, Any]:
        from PIL import Image
        p=Path(path)
        with Image.open(p) as im:
            im.verify()
        with Image.open(p) as im:
            return {"format":im.format,"width":im.width,"height":im.height,"mode":im.mode,"sha256":_sha_bytes(p.read_bytes())}

    @staticmethod
    def synthesize_speech(text: str, wav_path: str | Path, voice: str = "en") -> dict[str, Any]:
        exe=shutil.which("espeak-ng") or shutil.which("espeak")
        if not exe: raise FrontierError("speech synthesizer unavailable")
        p=Path(wav_path); p.parent.mkdir(parents=True,exist_ok=True)
        subprocess.run([exe,"-v",voice,"-w",str(p),text],check=True,capture_output=True)
        with wave.open(str(p),"rb") as w:
            duration=w.getnframes()/float(w.getframerate())
        if duration<=0: raise FrontierError("speech synthesis produced empty audio")
        return {"type":"speech-wav","duration_seconds":duration,"sha256":_sha_bytes(p.read_bytes())}

    @staticmethod
    def inspect_audio(wav_path: str | Path) -> dict[str, Any]:
        p=Path(wav_path)
        with wave.open(str(p),"rb") as w:
            frames=w.readframes(w.getnframes()); n=w.getnframes(); rate=w.getframerate(); channels=w.getnchannels(); width=w.getsampwidth()
        return {"frames":n,"rate":rate,"channels":channels,"sample_width":width,"duration_seconds":n/float(rate),"payload_sha256":_sha_bytes(frames)}

    @staticmethod
    def transcribe_vosk(wav_path: str | Path, model_path: str | Path) -> dict[str, Any]:
        from vosk import Model, KaldiRecognizer
        p=Path(wav_path)
        with wave.open(str(p),"rb") as w:
            if w.getnchannels()!=1 or w.getsampwidth()!=2:
                raise ValueError("Vosk adapter requires mono 16-bit PCM WAV")
            rec=KaldiRecognizer(Model(str(model_path)),w.getframerate())
            parts=[]
            while True:
                data=w.readframes(4000)
                if not data: break
                if rec.AcceptWaveform(data):
                    parts.append(json.loads(rec.Result()).get("text", ""))
            parts.append(json.loads(rec.FinalResult()).get("text", ""))
        text=" ".join(x.strip() for x in parts if x.strip()).strip()
        if not text: raise FrontierError("speech transcription produced no text")
        return {"text":text,"word_count":len(text.split()),"text_sha256":_sha_bytes(text.encode())}

    @staticmethod
    def probe_video(path: str | Path) -> dict[str, Any]:
        exe=shutil.which("ffprobe")
        if not exe: raise FrontierError("ffprobe unavailable")
        p=Path(path)
        proc=subprocess.run([exe,"-v","error","-show_entries","format=duration:stream=codec_type,codec_name,width,height","-of","json",str(p)],capture_output=True,text=True,check=True)
        meta=json.loads(proc.stdout)
        if float((meta.get("format") or {}).get("duration") or 0)<=0:
            raise FrontierError("video duration invalid")
        meta["sha256"]=_sha_bytes(p.read_bytes())
        return meta

    @staticmethod
    def camera_stream_frames(path_or_device: str | int, max_frames: int = 12) -> dict[str, Any]:
        import cv2
        cap=cv2.VideoCapture(path_or_device)
        if not cap.isOpened(): raise FrontierError("camera/video stream unavailable")
        hashes=[]; count=0
        try:
            while count<max_frames:
                ok,frame=cap.read()
                if not ok: break
                ok2,enc=cv2.imencode(".png",frame)
                if not ok2: raise FrontierError("frame encoding failed")
                hashes.append(_sha_bytes(bytes(enc))); count+=1
        finally:
            cap.release()
        if not hashes: raise FrontierError("camera/video stream produced no frames")
        return {"frames":count,"frame_hashes":hashes,"stream_sha256":_sha(hashes)}


# ---------------------------------------------------------------------------
# MUSITU Axiom live MCP adapter + specialist adapters
# ---------------------------------------------------------------------------

class AxiomMCPAdapter:
    def __init__(self, endpoint: str, access_token: str) -> None:
        p=urllib.parse.urlparse(endpoint)
        if p.scheme!="https" or not access_token:
            raise ValueError("HTTPS endpoint and access token required")
        self.endpoint=endpoint
        self.access_token=access_token

    def call_tool(self, name: str, arguments: Mapping[str, Any], request_id: str) -> dict[str, Any]:
        payload={"jsonrpc":"2.0","id":1,"method":"tools/call","params":{"name":name,"arguments":dict(arguments)}}
        req=urllib.request.Request(self.endpoint,data=_canonical(payload).encode(),headers={
            "Content-Type":"application/json","Accept":"application/json","Authorization":"Bearer "+self.access_token,
            "x-musitu-request-id":request_id,"User-Agent":"MUSITU-Axiom-Frontier-Specialist/5.0"},method="POST")
        try:
            with urllib.request.urlopen(req,timeout=45) as r:
                body=json.loads(r.read() or b"{}")
        except Exception as exc:
            raise FrontierError("Axiom MCP call failed") from exc
        result=body.get("result") or {}
        if result.get("isError") is True or body.get("error"):
            raise FrontierError("Axiom MCP returned error")
        return {"body":body,"body_sha256":_sha(body),"request_id":request_id}

    def execute(self, operation: str, args: Mapping[str, Any], request_id: str) -> dict[str, Any]:
        return self.call_tool("musitu_axiom_execute", {"operation":operation,"args":dict(args)}, request_id)


@dataclass(frozen=True)
class SpecialistToolBinding:
    specialist: str
    operation: str
    domain: str
    verified: bool


class SpecialistToolAdapter:
    def __init__(self, mcp: AxiomMCPAdapter, bindings: Sequence[SpecialistToolBinding]) -> None:
        self.mcp=mcp
        self.bindings={b.specialist:b for b in bindings}
        if len(self.bindings)!=len(bindings): raise ValueError("duplicate specialist binding")

    def run(self, specialist: str, args: Mapping[str, Any], request_id: str) -> dict[str, Any]:
        b=self.bindings.get(specialist)
        if b is None: raise FrontierError("specialist has no tool binding")
        if not b.verified: raise AuthorizationError("specialist tool binding is not verified")
        out=self.mcp.execute(b.operation,args,request_id)
        out.update({"specialist":specialist,"operation":b.operation,"domain":b.domain})
        return out


# ---------------------------------------------------------------------------
# Domain calibration + historical validation for digital twins
# ---------------------------------------------------------------------------

def _ols(x: Sequence[float], y: Sequence[float]) -> tuple[float,float,float]:
    if len(x)!=len(y) or len(x)<3: raise ValueError("at least three aligned observations required")
    mx=statistics.fmean(x); my=statistics.fmean(y)
    den=sum((a-mx)**2 for a in x)
    if den==0: raise ValueError("zero-variance regressor")
    slope=sum((a-mx)*(b-my) for a,b in zip(x,y))/den
    intercept=my-slope*mx
    pred=[intercept+slope*a for a in x]
    ss_res=sum((b-p)**2 for b,p in zip(y,pred)); ss_tot=sum((b-my)**2 for b in y)
    r2=1.0-ss_res/ss_tot if ss_tot else 1.0
    return intercept,slope,r2


class DomainTwinCalibrator:
    @staticmethod
    def company(twin_id: str, revenue: Sequence[float], operating_income: Sequence[float]) -> dict[str, Any]:
        intercept,slope,r2=_ols([float(x) for x in revenue],[float(x) for x in operating_income])
        model=CausalModel([LinearEquation("operating_income",intercept,{"revenue":slope})])
        twin=DigitalTwin(twin_id,"company",{"revenue":float(revenue[-1])},model,{"r2":r2})
        return {"twin":twin,"calibration":{"intercept":intercept,"slope":slope,"r2":r2,"n":len(revenue)}}

    @staticmethod
    def portfolio(twin_id: str, asset_returns: Mapping[str, Sequence[float]], weights: Mapping[str,float]) -> dict[str, Any]:
        names=sorted(asset_returns)
        if set(names)!=set(weights): raise ValueError("weights must match assets")
        n={len(asset_returns[k]) for k in names}
        if len(n)!=1 or next(iter(n))<3: raise ValueError("aligned return history required")
        wsum=sum(float(weights[k]) for k in names)
        if abs(wsum-1.0)>1e-6: raise ValueError("portfolio weights must sum to 1")
        portfolio=[sum(float(weights[k])*float(asset_returns[k][i]) for k in names) for i in range(next(iter(n)))]
        mean=statistics.fmean(portfolio); vol=statistics.stdev(portfolio)
        return {"twin_id":twin_id,"twin_type":"portfolio","mean_return":mean,"volatility":vol,"observations":len(portfolio),"assets":names,"history_sha256":_sha(asset_returns)}

    @staticmethod
    def economy(twin_id: str, driver: Sequence[float], outcome: Sequence[float], driver_name: str="driver", outcome_name: str="outcome") -> dict[str, Any]:
        intercept,slope,r2=_ols([float(x) for x in driver],[float(x) for x in outcome])
        model=CausalModel([LinearEquation(outcome_name,intercept,{driver_name:slope})])
        twin=DigitalTwin(twin_id,"economy",{driver_name:float(driver[-1])},model,{"r2":r2})
        return {"twin":twin,"calibration":{"intercept":intercept,"slope":slope,"r2":r2,"n":len(driver)}}

    @staticmethod
    def holdout_validate(train_x: Sequence[float], train_y: Sequence[float], test_x: Sequence[float], test_y: Sequence[float]) -> dict[str, Any]:
        intercept,slope,r2=_ols([float(x) for x in train_x],[float(x) for x in train_y])
        if len(test_x)!=len(test_y) or not test_x: raise ValueError("aligned holdout required")
        pred=[intercept+slope*float(x) for x in test_x]
        mae=statistics.fmean(abs(float(y)-p) for y,p in zip(test_y,pred))
        scale=statistics.fmean(abs(float(y)) for y in test_y) or 1.0
        return {"train_r2":r2,"holdout_mae":mae,"normalized_mae":mae/scale,"prediction_sha256":_sha(pred)}


# ---------------------------------------------------------------------------
# Sealed/unseen/adversarial/longitudinal comparative evaluation
# ---------------------------------------------------------------------------

@dataclass(frozen=True)
class ComparativeCase:
    case_id: str
    input: Mapping[str, Any]
    expected: Any
    tags: tuple[str,...]=()


class FrontierEvaluationHarness:
    @staticmethod
    def sealed_split(cases: Sequence[ComparativeCase], salt: str, holdout_fraction: float=.25) -> tuple[list[ComparativeCase],list[ComparativeCase]]:
        if not cases or not 0<holdout_fraction<1: raise ValueError("cases and valid holdout fraction required")
        ordered=sorted(cases,key=lambda c:hashlib.sha256(f"{salt}:{c.case_id}".encode()).hexdigest())
        n=max(1,min(len(ordered)-1,round(len(ordered)*holdout_fraction)))
        return ordered[n:],ordered[:n]

    @staticmethod
    def exact_score(actual: Any, expected: Any) -> float:
        return 1.0 if _canonical(actual)==_canonical(expected) else 0.0

    @staticmethod
    def numeric_score(actual: float, expected: float, atol: float=1e-9, rtol: float=1e-9) -> float:
        a=float(actual); e=float(expected); tol=atol+rtol*abs(e)
        return 1.0 if abs(a-e)<=tol else max(0.0,1.0-abs(a-e)/(abs(e)+1.0))

    @staticmethod
    def compare(cases: Sequence[ComparativeCase], runners: Mapping[str, Any], scorer: Any) -> dict[str, Any]:
        if len(runners)<2: raise ValueError("at least two independent runners required")
        rows=[]; means={name:[] for name in runners}
        for case in cases:
            row={"case_id":case.case_id,"scores":{},"outputs":{}}
            for name,runner in sorted(runners.items()):
                actual=runner(case.input); score=float(scorer(actual,case.expected))
                if not math.isfinite(score): raise ValueError("non-finite score")
                row["scores"][name]=score; row["outputs"][name]=_sha(actual); means[name].append(score)
            rows.append(row)
        summary={name:statistics.fmean(vals) for name,vals in means.items()}
        return {"count":len(rows),"means":summary,"rows":rows,"suite_sha256":_sha(rows)}

    @staticmethod
    def adversarial_gate(results: Sequence[Mapping[str,Any]]) -> dict[str,Any]:
        failed=[i for i,r in enumerate(results) if not bool(r.get("pass"))]
        return {"status":"PASS" if not failed else "FAIL","failed":failed,"count":len(results)}

    @staticmethod
    def longitudinal_gate(history: Sequence[Mapping[str,float]], metric: str, tolerance: float=0.0) -> dict[str,Any]:
        if len(history)<2: raise ValueError("at least two longitudinal points required")
        vals=[float(x[metric]) for x in history]
        regressions=[i for i in range(1,len(vals)) if vals[i]+tolerance<vals[i-1]]
        return {"status":"PASS" if not regressions else "FAIL","values":vals,"regressions":regressions}


__all__=[
    "AxiomMCPAdapter","ArtifactWorkbench","BrowserResult","CloudflareD1ProductionStore","ComparativeCase",
    "DomainTwinCalibrator","FrontierEvaluationHarness","LiveResearchAdapter","MultimodalWorkbench",
    "PlaywrightBrowserAdapter","RetrievalSnapshot","RetrievedContentFirewall","SandboxedComputerAdapter",
    "SQLiteProductionStore","SpecialistToolAdapter","SpecialistToolBinding",
]
